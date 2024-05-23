import cv2
import os
import argparse
import numpy as np
from PIL import Image
import imagehash
from pptx import Presentation
from pptx.util import Inches

def extract_frames(video_path, output_folder, interval_sec=1.5):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    cap = cv2.VideoCapture(video_path)
    fps = cap.get(cv2.CAP_PROP_FPS)
    interval = int(fps * interval_sec)

    frame_count = 0
    saved_count = 0

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break

        if frame_count % interval == 0:
            frame_filename = os.path.join(output_folder, f"frame_{saved_count}.jpg")
            cv2.imwrite(frame_filename, frame)
            saved_count += 1

        frame_count += 1

    cap.release()

def find_similar_images_and_rename(folder, threshold=5):
    image_hashes = []
    files_to_keep = []

    for filename in os.listdir(folder):
        if filename.endswith(".jpg"):
            filepath = os.path.join(folder, filename)
            image = Image.open(filepath)
            image_hash = imagehash.average_hash(image)

            is_similar = False
            for existing_hash in image_hashes:
                if image_hash - existing_hash <= threshold:
                    is_similar = True
                    break

            if not is_similar:
                image_hashes.append(image_hash)
                files_to_keep.append(filepath)

    for filename in os.listdir(folder):
        if filename.endswith(".jpg"):
            filepath = os.path.join(folder, filename)
            if filepath not in files_to_keep:
                os.remove(filepath)

    remaining_files = sorted([f for f in os.listdir(folder) if f.endswith(".jpg") and '_' in f and f.split('_')[1].split('.')[0].isdigit()], key=lambda x: int(x.split('_')[1].split('.')[0]))
    for i, filename in enumerate(remaining_files, start=1):
        file_extension = os.path.splitext(filename)[1]
        new_filename = f"{i}{file_extension}"
        old_filepath = os.path.join(folder, filename)
        new_filepath = os.path.join(folder, new_filename)
        os.rename(old_filepath, new_filepath)

def extract_ppt_from_image(image, output_path):
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edged = cv2.Canny(blurred, 50, 150)
    contours, _ = cv2.findContours(edged, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    max_rect = None
    max_area = 0

    for contour in contours:
        (x, y, w, h) = cv2.boundingRect(contour)
        area = w * h

        if area > max_area:
            max_area = area
            max_rect = (x, y, w, h)

    if max_rect is not None:
        x, y, w, h = max_rect
        ppt_region = image[y:y+h, x:x+w]
        cv2.imwrite(output_path, ppt_region)
        print(f"PPT区域已保存到: {output_path}")
    else:
        print("未找到PPT区域")

def process_images_in_folder(input_folder, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for filename in os.listdir(input_folder):
        if filename.endswith(".jpg") or filename.endswith(".png"):
            image_path = os.path.join(input_folder, filename)
            image = cv2.imread(image_path)

            if image is not None:
                output_path = os.path.join(output_folder, filename)
                extract_ppt_from_image(image, output_path)
            else:
                print(f"无法读取图片: {filename}")

def find_largest_image(input_folder, output_path):
    max_area = 0
    largest_image = None
    largest_image_path = None

    for filename in os.listdir(input_folder):
        if filename.endswith(".jpg") or filename.endswith(".png"):
            image_path = os.path.join(input_folder, filename)
            image = cv2.imread(image_path)

            if image is not None:
                height, width, _ = image.shape
                area = width * height

                if area > max_area:
                    max_area = area
                    largest_image = image
                    largest_image_path = image_path

    if largest_image is not None:
        cv2.imwrite(output_path, largest_image)
        print(f"最大尺寸的图片已保存到: {output_path}")
        print(f"最大尺寸的图片路径是: {largest_image_path}")
    else:
        print("未找到任何图片")

def find_template_in_image(image, template):
    result = cv2.matchTemplate(image, template, cv2.TM_CCOEFF_NORMED)
    min_val, max_val, min_loc, max_loc = cv2.minMaxLoc(result)
    top_left = max_loc
    height, width, _ = template.shape
    bottom_right = (top_left[0] + width, top_left[1] + height)
    return top_left, bottom_right

def crop_images_to_template(input_folder, template_path, output_folder):
    template = cv2.imread(template_path)
    if template is None:
        print("无法读取模板图片")
        return

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for filename in os.listdir(input_folder):
        if filename.endswith(".jpg") or filename.endswith(".png"):
            image_path = os.path.join(input_folder, filename)
            image = cv2.imread(image_path)

            if image is not None:
                top_left, bottom_right = find_template_in_image(image, template)
                cropped_image = image[top_left[1]:bottom_right[1], top_left[0]:bottom_right[0]]
                output_path = os.path.join(output_folder, filename)
                cv2.imwrite(output_path, cropped_image)
                print(f"{filename} 已裁剪并保存到 {output_path}")
            else:
                print(f"无法读取图片: {filename}")

def images_to_ppt(image_folder, output_ppt):
    presentation = Presentation()
    image_files = sorted([f for f in os.listdir(image_folder) if f.endswith((".jpg", ".png", ".jpeg"))], key=lambda x: int(os.path.splitext(x)[0]))

    for filename in image_files:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        img_path = os.path.join(image_folder, filename)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

    presentation.save(output_ppt)
    print("PPT已保存至:", output_ppt)

def main(video_path, output_ppt):
    output_folder = "frames"
    largest_images_folder = "largest-images"
    cropped_folder = "Cropped-frames"
    largest_image_path = "output_largest_image.jpg"

    extract_frames(video_path, output_folder)
    find_similar_images_and_rename(output_folder, threshold=5)
    process_images_in_folder(output_folder, largest_images_folder)
    find_largest_image(largest_images_folder, largest_image_path)
    crop_images_to_template(output_folder, largest_image_path, cropped_folder)
    images_to_ppt(cropped_folder, output_ppt)

    print("完成！")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="从视频中提取帧，裁剪并生成PPT")
    parser.add_argument("video_path", type=str, help="输入视频的路径")
    parser.add_argument("output_ppt", type=str, help="输出PPT的路径")

    args = parser.parse_args()
    main(args.video_path, args.output_ppt)
