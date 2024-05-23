from pptx import Presentation
from pptx.util import Inches
import os
import argparse


def images_to_ppt(image_folder, output_ppt):
    # 创建一个PPT文件
    presentation = Presentation()

    # 获取文件夹中的图片，并按照文件名中的数字部分排序
    image_files = sorted([f for f in os.listdir(image_folder) if
                          f.endswith((".jpg", ".png", ".jpeg"))],
                         key=lambda x: int(os.path.splitext(x)[0]))

    # 遍历文件夹中的图片，将每张图片按顺序添加到PPT中
    for filename in image_files:
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[5])  # 使用空白幻灯片布局
        img_path = os.path.join(image_folder, filename)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                 width=Inches(10),
                                 height=Inches(7.5))  # 调整图片大小以适应幻灯片

    # 保存PPT文件
    presentation.save(output_ppt)
    print("PPT已保存至:", output_ppt)


def main(image_folder, output_ppt):
    images_to_ppt(image_folder, output_ppt)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="将图片文件夹中的图片生成PPT")
    parser.add_argument("image_folder", type=str, help="图片文件夹的路径")
    parser.add_argument("output_ppt", type=str, help="输出PPT的路径")

    args = parser.parse_args()
    main(args.image_folder, args.output_ppt)
