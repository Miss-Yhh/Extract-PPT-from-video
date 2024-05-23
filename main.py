import cv2
import os
import argparse
from PIL import Image
import imagehash
from pptx import Presentation
from pptx.util import Inches


def extract_frames(video_path, output_folder, interval_sec=1.5):
    # 创建一个文件夹来保存提取的帧
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # 打开视频文件
    cap = cv2.VideoCapture(video_path)

    # 获取视频帧率
    fps = cap.get(cv2.CAP_PROP_FPS)
    interval = int(fps * interval_sec)  # 每隔interval_sec秒提取一帧

    frame_count = 0
    saved_count = 0

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break

        if frame_count % interval == 0:
            frame_filename = os.path.join(output_folder,
                                          f"frame_{saved_count}.jpg")
            cv2.imwrite(frame_filename, frame)
            saved_count += 1

        frame_count += 1

    cap.release()


def find_similar_images_and_rename(folder, threshold=5):
    image_hashes = []
    files_to_keep = []

    # 计算每个图片的哈希值
    for filename in os.listdir(folder):
        if filename.endswith(".jpg"):
            filepath = os.path.join(folder, filename)
            image = Image.open(filepath)
            image_hash = imagehash.average_hash(image)

            is_similar = False
            for existing_hash in image_hashes:
                if image_hash - existing_hash <= threshold:
                    is_similar = True
                    break

            if not is_similar:
                image_hashes.append(image_hash)
                files_to_keep.append(filepath)

    # 删除不保留的图片
    for filename in os.listdir(folder):
        if filename.endswith(".jpg"):
            filepath = os.path.join(folder, filename)
            if filepath not in files_to_keep:
                os.remove(filepath)

    # 重新命名图片文件
    remaining_files = sorted(
        [f for f in os.listdir(folder) if f.endswith(".jpg") and '_' in f],
        key=lambda x: int(x.split('_')[1].split('.')[0]))
    for i, filename in enumerate(remaining_files, start=1):
        file_extension = os.path.splitext(filename)[1]
        new_filename = f"{i}{file_extension}"
        old_filepath = os.path.join(folder, filename)
        new_filepath = os.path.join(folder, new_filename)
        os.rename(old_filepath, new_filepath)


def images_to_ppt(image_folder, output_ppt):
    # 创建一个PPT文件
    presentation = Presentation()

    # 获取文件夹中的图片，并按照文件名中的数字部分排序
    image_files = sorted([f for f in os.listdir(image_folder) if
                          f.endswith((".jpg", ".png", ".jpeg")) and
                          f.split('.')[0].isdigit()],
                         key=lambda x: int(os.path.splitext(x)[0]))

    # 遍历文件夹中的图片，将每张图片按顺序添加到PPT中
    for filename in image_files:
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[5])  # 使用空白幻灯片布局
        img_path = os.path.join(image_folder, filename)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                 width=Inches(10),
                                 height=Inches(7.5))  # 调整图片大小以适应幻灯片

    # 保存PPT文件
    presentation.save(output_ppt)
    print("PPT已保存至:", output_ppt)


def main(video_path, output_ppt):
    output_folder = "frames"

    # 提取帧
    extract_frames(video_path, output_folder)

    # 删除相似图片并重命名剩余图片
    find_similar_images_and_rename(output_folder, threshold=5)

    # 将图片生成PPT
    images_to_ppt(output_folder, output_ppt)

    print("完成！")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="从视频中提取帧并生成PPT")
    parser.add_argument("video_path", type=str, help="输入视频的路径")
    parser.add_argument("output_ppt", type=str, help="输出PPT的路径")

    args = parser.parse_args()
    main(args.video_path, args.output_ppt)
